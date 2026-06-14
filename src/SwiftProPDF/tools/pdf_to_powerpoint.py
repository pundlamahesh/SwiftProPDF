from pathlib import Path
import shutil
import tempfile

from pptx import Presentation
from pptx.util import Inches

from SwiftProPDF.tools.exceptions import PdfConversionError
from SwiftProPDF.tools.pdf_to_images import pdf_to_images


def pdf_to_powerpoint(input_path: Path, output_path: Path, overwrite: bool = False) -> None:
    """Convert PDF pages to PowerPoint slides."""
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfConversionError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfConversionError(f"Output file already exists: {output_path}")

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        temp_dir = Path(tempfile.mkdtemp(prefix="pdf2ppt-"))

        try:
            image_paths = pdf_to_images(input_path, temp_dir, dpi=150)
            prs = Presentation()

            for img_path in image_paths:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(str(img_path), 0, 0, width=Inches(10), height=Inches(7.5))

            prs.save(str(output_path))

        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    except Exception as exc:
        raise PdfConversionError(f"Could not convert PDF to PowerPoint: {str(exc)}") from exc
