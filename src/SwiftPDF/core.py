from pathlib import Path
import shutil
import subprocess
import tempfile
import logging

from pypdf.errors import DependencyError
from pypdf import PdfReader, PdfWriter
from PIL import Image
import fitz  # PyMuPDF
from pdf2docx import Converter as PDFToWordConverter
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

logger = logging.getLogger(__name__)


# ============================================================================
# EXCEPTION CLASSES
# ============================================================================

class PdfUnlockError(Exception):
    """Raised when a PDF cannot be unlocked."""


class PdfLockError(Exception):
    """Raised when a PDF cannot be locked."""


class PdfSplitError(Exception):
    """Raised when a PDF cannot be split."""


class PdfMergeError(Exception):
    """Raised when PDFs cannot be merged."""


class PdfCompressError(Exception):
    """Raised when a PDF cannot be compressed."""


class PdfConversionError(Exception):
    """Raised when PDF conversion fails."""


class ImageConversionError(Exception):
    """Raised when image conversion fails."""


class OfficeConversionError(Exception):
    """Raised when Office document conversion fails."""


class PdfEditError(Exception):
    """Raised when PDF editing fails."""


class QrCodeError(Exception):
    """Raised when QR code generation fails."""


# ============================================================================
# PDF SPLIT
# ============================================================================

def parse_page_ranges(page_ranges: str, page_count: int) -> list[int]:
    """Parse 1-based page ranges into unique 0-based page indexes."""
    if page_count < 1:
        raise PdfSplitError("PDF does not contain any pages.")

    if not page_ranges.strip():
        raise PdfSplitError("Enter page ranges such as 1-3,5.")

    selected_pages: list[int] = []
    seen_pages: set[int] = set()

    for part in page_ranges.split(","):
        part = part.strip()
        if not part:
            continue

        if "-" in part:
            start_text, end_text = [value.strip() for value in part.split("-", 1)]
            if not start_text.isdigit() or not end_text.isdigit():
                raise PdfSplitError("Page ranges must use numbers, for example 1-3,5.")

            start_page = int(start_text)
            end_page = int(end_text)
            if start_page > end_page:
                raise PdfSplitError("Page range start cannot be greater than the end.")

            pages = range(start_page, end_page + 1)
        else:
            if not part.isdigit():
                raise PdfSplitError("Page ranges must use numbers, for example 1-3,5.")
            pages = range(int(part), int(part) + 1)

        for page_number in pages:
            if page_number < 1 or page_number > page_count:
                raise PdfSplitError(f"Page {page_number} is outside this PDF's 1-{page_count} page range.")

            page_index = page_number - 1
            if page_index not in seen_pages:
                selected_pages.append(page_index)
                seen_pages.add(page_index)

    if not selected_pages:
        raise PdfSplitError("Enter at least one page number.")

    return selected_pages
# ============================================================================
# PDF UNLOCK
# ============================================================================

def unlock_pdf(input_path: Path, output_path: Path, password: str, overwrite: bool = False) -> None:
    """Unlock a PDF with the given password and write an unencrypted copy."""
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfUnlockError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfUnlockError(f"Output file already exists: {output_path}")

    try:
        reader = PdfReader(str(input_path))
    except Exception as exc:
        raise PdfUnlockError(f"Could not read PDF: {input_path}") from exc

    try:
        if reader.is_encrypted:
            decrypt_result = reader.decrypt(password)
            if decrypt_result == 0:
                raise PdfUnlockError("Could not unlock PDF. Check the password.")
        elif password:
            # Accept an unnecessary password so scripts can handle mixed encrypted/plain PDFs.
            pass

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
    except DependencyError as exc:
        raise PdfUnlockError(
            "This PDF uses AES encryption. Install crypto support with: "
            'python -m pip install --upgrade -e . or python -m pip install "pypdf[crypto]" cryptography'
        ) from exc

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        with output_path.open("wb") as output_file:
            writer.write(output_file)
    except Exception as exc:
        raise PdfUnlockError(f"Could not write unlocked PDF: {output_path}") from exc


# ============================================================================
# PDF LOCK
# ============================================================================

def lock_pdf(input_path: Path, output_path: Path, password: str, overwrite: bool = False) -> None:
    """Encrypt a PDF with the given password."""
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfLockError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfLockError(f"Output file already exists: {output_path}")

    if not password:
        raise PdfLockError("Enter a password to lock the PDF.")

    try:
        reader = PdfReader(str(input_path))
    except Exception as exc:
        raise PdfLockError(f"Could not read PDF: {input_path}") from exc

    if reader.is_encrypted:
        raise PdfLockError("This PDF is already encrypted. Unlock it before locking again.")

    try:
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(password)
    except Exception as exc:
        raise PdfLockError("Could not lock this PDF.") from exc

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        with output_path.open("wb") as output_file:
            writer.write(output_file)
    except Exception as exc:
        raise PdfLockError(f"Could not write locked PDF: {output_path}") from exc
# ============================================================================
# PDF SPLIT
# ============================================================================

def split_pdf(input_path: Path, output_path: Path, page_ranges: str, overwrite: bool = False) -> None:
    """Write a new PDF containing only the requested pages."""
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfSplitError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfSplitError(f"Output file already exists: {output_path}")

    try:
        reader = PdfReader(str(input_path))
    except Exception as exc:
        raise PdfSplitError(f"Could not read PDF: {input_path}") from exc

    if reader.is_encrypted:
        raise PdfSplitError("Encrypted PDFs must be unlocked before splitting.")

    try:
        selected_pages = parse_page_ranges(page_ranges, len(reader.pages))
        writer = PdfWriter()
        for page_index in selected_pages:
            writer.add_page(reader.pages[page_index])
    except PdfSplitError:
        raise
    except Exception as exc:
        raise PdfSplitError("Could not split this PDF.") from exc

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        with output_path.open("wb") as output_file:
            writer.write(output_file)
    except Exception as exc:
        raise PdfSplitError(f"Could not write split PDF: {output_path}") from exc


# ============================================================================
# MERGE PDF
# ============================================================================

def merge_pdfs(input_paths: list[Path], output_path: Path, overwrite: bool = False) -> None:
    """Merge multiple PDFs into a single PDF."""
    input_paths = [Path(p) for p in input_paths]
    output_path = Path(output_path)

    if output_path.exists() and not overwrite:
        raise PdfMergeError(f"Output file already exists: {output_path}")

    if not input_paths:
        raise PdfMergeError("No PDFs to merge.")

    try:
        writer = PdfWriter()
        
        for input_path in input_paths:
            if not input_path.exists():
                raise PdfMergeError(f"Input PDF does not exist: {input_path}")
            
            try:
                reader = PdfReader(str(input_path))
                
                if reader.is_encrypted:
                    raise PdfMergeError(f"PDF {input_path.name} is encrypted. Unlock it first.")
                
                for page in reader.pages:
                    writer.add_page(page)
            except PdfMergeError:
                raise
            except Exception as exc:
                raise PdfMergeError(f"Could not read PDF: {input_path}") from exc
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with output_path.open("wb") as output_file:
            writer.write(output_file)
            
    except PdfMergeError:
        raise
    except Exception as exc:
        raise PdfMergeError("Could not merge PDFs.") from exc


# ============================================================================
# COMPRESS PDF
# ============================================================================
def compress_pdf(input_path: Path, output_path: Path, level: str = "medium", overwrite: bool = False) -> None:
    """Compress PDF.

    Levels:
        low    -> Fast lossless optimization (PyMuPDF)
        medium -> Ghostscript balanced compression
        high   -> Ghostscript aggressive compression
    """

    import shutil
    import subprocess

    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfCompressError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfCompressError(f"Output file already exists: {output_path}")

    if level not in ("low", "medium", "high"):
        raise PdfCompressError("Compression level must be: low, medium, or high")

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # LOW = PyMuPDF optimization
        if level == "low":
            with fitz.open(str(input_path)) as doc:
                if doc.needs_pass:
                    raise PdfCompressError(
                        "Encrypted PDFs must be unlocked before compressing."
                    )

                doc.save(
                    str(output_path),
                    garbage=4,
                    clean=True,
                    deflate=True,
                    deflate_images=True,
                    deflate_fonts=True,
                )
            return

        # MEDIUM/HIGH = Ghostscript
        gs_binary = (
            shutil.which("gs")
            or shutil.which("gswin64c")
            or shutil.which("gswin32c")
            )

        if not gs_binary:
            raise PdfCompressError(
                "Ghostscript is not installed on the server."
            )

        pdf_setting = {
            "medium": "/ebook",
            "high": "/screen",
        }[level]

        subprocess.run(
            [
                gs_binary,
                "-sDEVICE=pdfwrite",
                "-dCompatibilityLevel=1.4",
                f"-dPDFSETTINGS={pdf_setting}",
                "-dNOPAUSE",
                "-dQUIET",
                "-dBATCH",
                f"-sOutputFile={output_path}",
                str(input_path),
            ],
            check=True,
            capture_output=True,
        )

        if not output_path.exists():
            raise PdfCompressError(
                "Compression completed but output file was not created."
            )

    except PdfCompressError:
        raise

    except subprocess.CalledProcessError as exc:
        error_msg = exc.stderr.decode(errors="ignore")
        raise PdfCompressError(
            f"Ghostscript compression failed: {error_msg}"
        ) from exc

    except Exception as exc:
        raise PdfCompressError(
            f"Could not compress PDF: {str(exc)}"
        ) from exc

# ============================================================================
# COMPRESS IMAGES
# ============================================================================

def compress_image(input_path: Path, output_path: Path, level: str = "medium", overwrite: bool = False) -> None:
    """Compress an image file.

    Levels:
        low    -> Mild compression, high quality
        medium -> Balanced compression
        high   -> Aggressive compression
    """

    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise ImageConversionError(f"Input image does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise ImageConversionError(f"Output file already exists: {output_path}")

    if level not in ("low", "medium", "high"):
        raise ImageConversionError("Compression level must be: low, medium, or high")

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)

        quality_map = {
            "low": 85,
            "medium": 70,
            "high": 50,
        }

        resize_map = {
            "low": 3000,
            "medium": 2000,
            "high": 1200,
        }

        compress_level_map = {
            "low": 3,
            "medium": 6,
            "high": 9,
        }

        with Image.open(input_path) as img:
            img_format = (img.format or input_path.suffix.lstrip(".")).upper()

            max_dimension = resize_map[level]

            # Resize large images
            if max(img.size) > max_dimension:
                img.thumbnail(
                    (max_dimension, max_dimension),
                    Image.Resampling.LANCZOS,
                )

            # JPEG
            if img_format in ("JPEG", "JPG"):
                rgb = img.convert("RGB")

                rgb.save(
                    output_path,
                    format="JPEG",
                    quality=quality_map[level],
                    optimize=True,
                    progressive=True,
                )
                return

            # PNG
            if img_format == "PNG":
                img.save(
                    output_path,
                    format="PNG",
                    optimize=True,
                    compress_level=compress_level_map[level],
                )
                return

            # WEBP
            if img_format == "WEBP":
                img.save(
                    output_path,
                    format="WEBP",
                    quality=quality_map[level],
                    optimize=True,
                    method=6,
                )
                return

            # GIF
            if img_format == "GIF":
                img.save(
                    output_path,
                    format="GIF",
                    optimize=True,
                )
                return

            # Fallback
            img.save(
                output_path,
                format=img_format,
                optimize=True,
            )

    except Exception as exc:
        raise ImageConversionError(
            f"Could not compress image: {str(exc)}"
        ) from exc

# ============================================================================
# PDF TO IMAGES (PDF → JPG)
# ============================================================================

def pdf_to_images(input_path: Path, output_dir: Path, dpi: int = 150, overwrite: bool = False) -> list[Path]:
    """Convert all PDF pages to JPG images.
    
    Returns list of output image paths.
    """
    input_path = Path(input_path)
    output_dir = Path(output_dir)

    if not input_path.exists():
        raise ImageConversionError(f"Input PDF does not exist: {input_path}")

    if dpi < 72 or dpi > 300:
        raise ImageConversionError("DPI must be between 72 and 300")

    try:
        output_paths = []
        doc = fitz.open(str(input_path))
        
        output_dir.mkdir(parents=True, exist_ok=True)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            
            output_path = output_dir / f"page_{page_num + 1:04d}.jpg"
            pix.save(str(output_path))
            output_paths.append(output_path)
        
        doc.close()
        return output_paths
        
    except Exception as exc:
        raise ImageConversionError(f"Could not convert PDF to images: {str(exc)}") from exc


# ============================================================================
# IMAGES TO PDF (JPG → PDF)
# ============================================================================

def images_to_pdf(input_paths: list[Path], output_path: Path, overwrite: bool = False) -> None:
    """Convert multiple JPG/PNG images to a single PDF."""
    input_paths = [Path(p) for p in input_paths]
    output_path = Path(output_path)

    if output_path.exists() and not overwrite:
        raise ImageConversionError(f"Output file already exists: {output_path}")

    if not input_paths:
        raise ImageConversionError("No images provided.")

    try:
        images = []
        
        for img_path in input_paths:
            if not img_path.exists():
                raise ImageConversionError(f"Image does not exist: {img_path}")
            
            if not img_path.suffix.lower() in (".jpg", ".jpeg", ".png", ".gif", ".bmp"):
                raise ImageConversionError(f"Unsupported image format: {img_path.suffix}")
            
            img = Image.open(str(img_path))
            
            # Convert to RGB if necessary
            if img.mode in ("RGBA", "LA", "P"):
                rgb_img = Image.new("RGB", img.size, (255, 255, 255))
                rgb_img.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
                images.append(rgb_img)
            else:
                images.append(img.convert("RGB"))
        
        if not images:
            raise ImageConversionError("No valid images to convert.")
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        images[0].save(str(output_path), save_all=True, append_images=images[1:], format="PDF")
        
    except ImageConversionError:
        raise
    except Exception as exc:
        raise ImageConversionError(f"Could not convert images to PDF: {str(exc)}") from exc
# ============================================================================
# QR CODE GENERATION
# ============================================================================

def generate_qr_code(data: str, output_path: Path, size: int = 500, overwrite: bool = False) -> None:
    """Generate a QR code PNG image from a URL or text."""
    output_path = Path(output_path)

    if output_path.exists() and not overwrite:
        raise QrCodeError(f"Output file already exists: {output_path}")

    try:
        import qrcode
        from qrcode.constants import ERROR_CORRECT_M
    except ImportError as exc:
        raise QrCodeError("Install the qrcode package to generate QR codes.") from exc

    try:
        qr = qrcode.QRCode(
            version=None,
            error_correction=ERROR_CORRECT_M,
            box_size=10,
            border=4,
        )
        qr.add_data(data)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")

        if hasattr(img, "convert"):
            img = img.convert("RGB")

        resample = getattr(getattr(Image, "Resampling", Image), "LANCZOS", Image.LANCZOS)
        img = img.resize((size, size), resample=resample)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        img.save(output_path, format="PNG")
    except Exception as exc:
        raise QrCodeError("Could not generate the QR code.") from exc


# ============================================================================
# PDF TO WORD
# ============================================================================

def pdf_to_word(input_path: Path, output_path: Path, overwrite: bool = False) -> None:
    """Convert PDF to Word document (.docx)."""
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfConversionError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfConversionError(f"Output file already exists: {output_path}")

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Use pdf2docx for conversion
        converter = PDFToWordConverter(str(input_path))
        converter.convert(str(output_path), start=0, end=None)
        converter.close()
        
    except Exception as exc:
        raise PdfConversionError(f"Could not convert PDF to Word: {str(exc)}") from exc


# ============================================================================
# PDF TO POWERPOINT
# ============================================================================

def pdf_to_powerpoint(input_path: Path, output_path: Path, overwrite: bool = False) -> None:
    """Convert PDF to PowerPoint presentation (.pptx).
    
    Each PDF page becomes a slide with the page image.
    """
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfConversionError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfConversionError(f"Output file already exists: {output_path}")

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Create temporary directory for images
        temp_dir = Path(tempfile.mkdtemp(prefix="pdf2ppt-"))
        
        try:
            # Convert PDF pages to images
            image_paths = pdf_to_images(input_path, temp_dir, dpi=150)
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            for img_path in image_paths:
                # Add blank slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add image to slide
                slide.shapes.add_picture(str(img_path), 0, 0, width=Inches(10), height=Inches(7.5))
            
            prs.save(str(output_path))
            
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)
        
    except Exception as exc:
        raise PdfConversionError(f"Could not convert PDF to PowerPoint: {str(exc)}") from exc


# ============================================================================
# PDF TO EXCEL
# ============================================================================

def pdf_to_excel(input_path: Path, output_path: Path, overwrite: bool = False) -> None:
    """Convert PDF content to Excel workbook (.xlsx) - one sheet per page."""
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfConversionError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfConversionError(f"Output file already exists: {output_path}")

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Open PDF with PyMuPDF
        doc = fitz.open(str(input_path))
        wb = Workbook()
        wb.remove(wb.active)  # Remove default sheet
        
        # Extract text from each page and create a sheet
        for page_idx, page in enumerate(doc):
            text = page.get_text()
            
            sheet_name = f"Page_{page_idx + 1}"
            if len(sheet_name) > 31:  # Excel sheet name limit
                sheet_name = sheet_name[:31]
            
            ws = wb.create_sheet(sheet_name)
            
            # Split text by lines and add to sheet
            lines = text.split('\n')
            for row_idx, line in enumerate(lines, 1):
                # Split by whitespace to create columns
                cells = line.split() if line.strip() else []
                for col_idx, cell_val in enumerate(cells, 1):
                    ws.cell(row=row_idx, column=col_idx, value=cell_val)
        
        doc.close()
        wb.save(str(output_path))
        
    except Exception as exc:
        raise PdfConversionError(f"Could not convert PDF to Excel: {str(exc)}") from exc


# ============================================================================
# OFFICE TO PDF (WORD, POWERPOINT, EXCEL)
# ============================================================================

def office_to_pdf(input_path: Path, output_path: Path, overwrite: bool = False) -> None:
    """Convert Office documents (DOCX, XLSX, PPTX) to PDF using LibreOffice."""
    import os
    import shutil
    import tempfile
    import subprocess
    from pathlib import Path

    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise OfficeConversionError(f"Input file does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise OfficeConversionError(f"Output file already exists: {output_path}")

    supported_formats = (".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt")

    if input_path.suffix.lower() not in supported_formats:
        raise OfficeConversionError(
            f"Unsupported format: {input_path.suffix}"
        )

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)

        temp_dir = Path(tempfile.mkdtemp(prefix="office2pdf-"))
        

        try:
            #print("PATH =", os.environ.get("PATH"))
            #print("which soffice =", shutil.which("soffice"))
            #print("which libreoffice =", shutil.which("libreoffice"))
            libreoffice_cmd = (
                shutil.which("soffice")
                or shutil.which("libreoffice")
            )

            if not libreoffice_cmd:
                raise OfficeConversionError(
                    "LibreOffice is not installed or not available in PATH."
                )

            result = subprocess.run(
                [
                    libreoffice_cmd,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(temp_dir),
                    str(input_path),
                ],
                capture_output=True,
                text=True,
                timeout=60,
            )

            if result.returncode != 0:
                raise OfficeConversionError(
                    f"LibreOffice conversion failed:\n"
                    f"STDOUT: {result.stdout}\n"
                    f"STDERR: {result.stderr}"
                )

            temp_pdf = temp_dir / f"{input_path.stem}.pdf"

            if not temp_pdf.exists():
                raise OfficeConversionError(
                    f"Conversion completed but PDF not found: {temp_pdf}"
                )

            shutil.move(str(temp_pdf), str(output_path))

        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    except OfficeConversionError:
        raise

    except subprocess.TimeoutExpired as exc:
        raise OfficeConversionError(
            "Conversion timeout. File may be too large."
        ) from exc

    except Exception as exc:
        raise OfficeConversionError(
            f"Could not convert Office document to PDF: {str(exc)}"
        ) from exc


# ============================================================================
# EDIT PDF (BASIC OPERATIONS)
# ============================================================================

def rotate_pdf_pages(input_path: Path, output_path: Path, page_ranges: str, angle: int, overwrite: bool = False) -> None:
    """Rotate PDF pages by specified angle (90, 180, 270, -90)."""
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfEditError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfEditError(f"Output file already exists: {output_path}")

    if angle not in (90, 180, 270, -90):
        raise PdfEditError("Rotation angle must be 90, 180, 270, or -90 degrees")

    try:
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        
        if reader.is_encrypted:
            raise PdfEditError("Cannot edit encrypted PDFs.")
        
        selected_pages = (
            list(range(len(reader.pages)))
            if not page_ranges.strip()
            else parse_page_ranges(page_ranges, len(reader.pages))
        )
        selected_set = set(selected_pages)
        
        for page_idx, page in enumerate(reader.pages):
            if page_idx in selected_set:
                page.rotate(angle)
            writer.add_page(page)
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with output_path.open("wb") as f:
            writer.write(f)
    
    except (PdfEditError, PdfSplitError) as exc:
        if isinstance(exc, PdfSplitError):
            raise PdfEditError(str(exc)) from exc
        raise
    except Exception as exc:
        raise PdfEditError(f"Could not rotate PDF pages: {str(exc)}") from exc

# ============================================================================
# DELETE PDF PAGES
# ============================================================================
def delete_pdf_pages(input_path: Path, output_path: Path, page_ranges: str, overwrite: bool = False) -> None:
    """Delete specified pages from PDF."""
    input_path = Path(input_path)
    output_path = Path(output_path)

    if not input_path.exists():
        raise PdfEditError(f"Input PDF does not exist: {input_path}")

    if output_path.exists() and not overwrite:
        raise PdfEditError(f"Output file already exists: {output_path}")

    try:
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        
        if reader.is_encrypted:
            raise PdfEditError("Cannot edit encrypted PDFs.")
        
        # Parse page ranges to delete
        pages_to_delete = set(parse_page_ranges(page_ranges, len(reader.pages)))
        
        # Add all pages except deleted ones
        for page_idx, page in enumerate(reader.pages):
            if page_idx not in pages_to_delete:
                writer.add_page(page)
        
        if len(writer.pages) == 0:
            raise PdfEditError("Cannot delete all pages from PDF.")
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with output_path.open("wb") as f:
            writer.write(f)
    
    except (PdfEditError, PdfSplitError) as exc:
        if isinstance(exc, PdfSplitError):
            raise PdfEditError(str(exc)) from exc
        raise
    except Exception as exc:
        raise PdfEditError(f"Could not delete PDF pages: {str(exc)}") from exc
